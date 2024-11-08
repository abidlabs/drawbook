"""
Core functionality for the drawbook library.
"""

from pathlib import Path
from typing import List, Literal
import tempfile
import io
import requests
import warnings
from tqdm import tqdm
from PIL import Image
import huggingface_hub
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import gradio as gr
from PIL import ImageDraw, ImageFont

class Book:
    """A class representing a children's book that can be exported to PowerPoint."""
    
    def __init__(
        self,
        title: str = "Untitled Book",
        pages: List[str] = None,
        title_illustration: str | Literal[False] | None = None,
        illustrations: List[str | None | Literal[False]] = None,
        lora: str = "SebastianBodza/Flux_Aquarell_Watercolor_v2",
        author: str | None = None
    ):
        """
        Initialize a new Book.
        
        Args:
            title: The book's title
            pages: List of strings containing text for each page
            illustrations: List of illustration paths or placeholders
                         (str for path, None for pending, False for no illustration)
            lora: The LoRA model on Hugging Face to use for illustrations
            author: The book's author name
        """
        self.title = title
        self.pages = pages or []
        self.illustrations = illustrations or []
        self.lora = lora
        self.title_illustration = title_illustration
        self.author = author
        
        # Ensure illustrations list matches pages length
        while len(self.illustrations) < len(self.pages):
            self.illustrations.append(None)

    def _get_prompt(self, text: str) -> str:
        """Get the prompt for the given text using the LoRA model."""
        if self.lora == "SebastianBodza/Flux_Aquarell_Watercolor_v2":
            return f"A AQUACOLTOK watercolor painting with a white background to illustrate the following text in a children's book: {text}"
        else:
            return f"An illustration of the following text in a children's book: {text}"
    
    def export(self, filename: str | Path | None = None) -> None:
        """
        Export the book to a PowerPoint file.
        
        Args:
            filename: Optional path where to save the file. If None, creates in temp directory.
        """
        if filename is None:
            # Create temp file with .pptx extension
            temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
            output_path = Path(temp_file.name)
            temp_file.close()
        else:
            # Convert to Path object and resolve to absolute path
            output_path = Path(filename).resolve()
            # Ensure parent directories exist
            output_path.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Add diagonal striped border at the top
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.5)
        )
        border.fill.solid()
        border.fill.fore_color.rgb = RGBColor(128, 0, 0)  # Maroon
        
        # Add title with adjusted positioning and z-order
        title = slide.shapes.title
        title.top = Inches(1.0)  # Increased spacing from top
        title.height = Inches(2.0)  # Increased height to accommodate two lines
        
        # Split title into first word and rest
        title_parts = self.title.split(maxsplit=1)
        first_word = title_parts[0]
        rest_of_title = title_parts[1] if len(title_parts) > 1 else ""
        
        # Add first word
        p1 = title.text_frame.paragraphs[0]
        p1.text = first_word
        p1.font.name = "Trebuchet MS"
        p1.font.color.rgb = RGBColor(128, 0, 0)
        p1.font.size = Inches(0.5)
        p1.alignment = PP_ALIGN.CENTER
        
        # Add rest of title
        if rest_of_title:
            p2 = title.text_frame.add_paragraph()
            p2.text = rest_of_title
            p2.font.name = "Trebuchet MS"
            p2.font.color.rgb = RGBColor(128, 0, 0)
            p2.font.size = Inches(0.5)
            p2.alignment = PP_ALIGN.CENTER
        
        # Add title illustration if available
        if isinstance(self.title_illustration, str):
            try:
                slide.shapes.add_picture(
                    self.title_illustration,
                    Inches(2), Inches(2.5),
                    Inches(6), Inches(5)     # Reduced height to make room for author
                )
            except Exception as e:
                print(f"Warning: Could not add title illustration: {e}")
        
        # Add author with adjusted positioning
        if self.author is not None:
            author_box = slide.shapes.add_textbox(
                Inches(0), Inches(7.5),  # Moved up from bottom
                Inches(10), Inches(0.5)
            )
            author_frame = author_box.text_frame
            author_frame.text = f"Written by {self.author}"
            author_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            author_frame.paragraphs[0].font.name = "Trebuchet MS"
            author_frame.paragraphs[0].font.size = Inches(0.2)
        
        # Add content slides
        content_slide_layout = prs.slide_layouts[5]  # Blank layout
        
        for page_num, (text, illustration) in enumerate(zip(self.pages, self.illustrations)):
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Determine if text should be above or below illustration
            text_on_top = page_num % 2 == 0
            text_y = Inches(0.5) if text_on_top else Inches(4.5)
            illustration_y = Inches(2) if text_on_top else Inches(0.5)
            
            # Split text into sentences and join with newlines
            sentences = text.replace('. ', '.\n').split('\n')
            
            # Special formatting for first page
            if page_num == 0 and text:
                # Split first character from first sentence
                first_char = sentences[0][0]
                first_sentence_rest = sentences[0][1:]
                
                p = slide.shapes.title.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = first_char
                run.font.size = Inches(0.3)  # Make first letter larger
                run.font.name = "Trebuchet MS"
                
                run = p.add_run()
                run.text = first_sentence_rest
                run.font.size = Inches(0.25)  # Regular text size
                run.font.name = "Trebuchet MS"
                
                # Add remaining sentences as new paragraphs
                for sentence in sentences[1:]:
                    p = slide.shapes.title.text_frame.add_paragraph()
                    p.text = sentence
                    p.font.name = "Trebuchet MS"
                    p.font.size = Inches(0.25)
                    p.alignment = PP_ALIGN.CENTER
            else:
                # Add each sentence as a separate paragraph
                first_paragraph = True
                for sentence in sentences:
                    if first_paragraph:
                        p = slide.shapes.title.text_frame.paragraphs[0]
                        first_paragraph = False
                    else:
                        p = slide.shapes.title.text_frame.add_paragraph()
                    p.text = sentence
                    p.font.name = "Trebuchet MS"
                    p.font.size = Inches(0.25)
                    p.alignment = PP_ALIGN.CENTER
            
            # Add illustration if available
            if isinstance(illustration, str):
                try:
                    slide.shapes.add_picture(
                        illustration,
                        Inches(1), illustration_y,
                        Inches(8), Inches(4)
                    )
                except Exception as e:
                    print(f"Warning: Could not add illustration on page {page_num + 1}: {e}")
            
            # Add page number at bottom center
            page_number = page_num + 1  # Add 1 since page_num is 0-based
            page_num_box = slide.shapes.add_textbox(
                Inches(0), Inches(6.5),  # Y position near bottom of slide
                Inches(10), Inches(0.5)  # Full width of slide for centering
            )
            page_num_frame = page_num_box.text_frame
            page_num_frame.text = str(page_number)
            page_num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            page_num_frame.paragraphs[0].font.name = "Trebuchet MS"
            page_num_frame.paragraphs[0].font.size = Inches(0.15)  # Slightly smaller than main text
        
        # Save the presentation
        prs.save(str(output_path))
        print(f"Book exported to: {output_path.absolute()}")
    
    def __len__(self) -> int:
        """Return the number of pages in the book."""
        return len(self.pages)

    def illustrate(self, save_dir: str | Path | None = None) -> None:
        """
        Generate illustrations for all pages using the Hugging Face Inference API.
        
        Args:
            save_dir: Optional directory to save the generated images. 
                     If None, creates a temporary directory.
        """
        # Get HF token
        token = huggingface_hub.get_token()
        if not token:
            warnings.warn("No Hugging Face token found. Please login using `huggingface-cli login` or set the HF_TOKEN environment variable. Otherwise, you may be rate limited.")

        API_URL = f"https://api-inference.huggingface.co/models/{self.lora}"
        headers = {"Authorization": f"Bearer {token}"}

        # Create save directory if provided
        if save_dir:
            save_dir = Path(save_dir)
            save_dir.mkdir(parents=True, exist_ok=True)
        else:
            save_dir = Path(tempfile.mkdtemp())

        print("Generating illustrations... This could take a few minutes.")

        # Create a list of tasks for the progress bar
        tasks = []
        if self.title_illustration is None:
            tasks.append(("title", self.title, None))
        tasks.extend((f"page_{i+1}", text, current_illust) 
                     for i, (text, current_illust) in enumerate(zip(self.pages, self.illustrations)))

        for task_name, text, current_illust in tqdm(tasks, desc="Generating illustrations"):
            # Skip if illustration already exists or is explicitly disabled
            if isinstance(current_illust, str) or current_illust is False:
                continue
                
            try:
                # Query the API
                prompt = self._get_prompt(text)
                response = requests.post(API_URL, headers=headers, json={"inputs": prompt})
                if response.status_code != 200:
                    print(f"Warning: Failed to generate illustration for {task_name}: {response.text}")
                    continue
                    
                # Save the image
                image = Image.open(io.BytesIO(response.content))
                image_path = save_dir / f"{task_name}.png"
                image.save(image_path)
                
                # Update the appropriate illustration reference
                if task_name == "title":
                    self.title_illustration = str(image_path)
                else:
                    page_num = int(task_name.split('_')[1]) - 1
                    self.illustrations[page_num] = str(image_path)
                    
            except Exception as e:
                print(f"Warning: Error generating illustration for {task_name}: {e}")
                continue

        print(f"Illustrations saved to: {save_dir}")

    def _create_page_image(self, page_num: int) -> Image.Image:
        """Create a PIL Image representation of a book page."""
        # Create a white background image (same aspect ratio as PPT slide)
        img = Image.new('RGB', (1000, 750), 'white')
        draw = ImageDraw.Draw(img)
        
        # Try to load a basic font
        try:
            font = ImageFont.truetype("Arial.ttf", 36)
            small_font = ImageFont.truetype("Arial.ttf", 24)
        except:
            font = ImageFont.load_default()
            small_font = ImageFont.load_default()

        text = self.pages[page_num]
        illustration = self.illustrations[page_num]
        
        # Determine text position (alternating top/bottom like in export())
        text_on_top = page_num % 2 == 0
        text_y = 50 if text_on_top else 450
        illustration_y = 200 if text_on_top else 50

        # Split text into sentences and calculate total height
        sentences = text.replace('. ', '.\n').split('\n')
        line_height = 40  # Adjust as needed
        
        # Draw each sentence on a new line
        for i, sentence in enumerate(sentences):
            text_bbox = draw.textbbox((0, 0), sentence, font=font)
            text_width = text_bbox[2] - text_bbox[0]
            text_x = (1000 - text_width) // 2  # Center text
            current_y = text_y + (i * line_height)
            draw.text((text_x, current_y), sentence, fill='black', font=font)

        # Add illustration if available
        if isinstance(illustration, str):
            try:
                illust = Image.open(illustration)
                illust = illust.resize((800, 400), Image.Resampling.LANCZOS)
                img.paste(illust, (100, illustration_y))
            except Exception as e:
                draw.text((100, illustration_y), f"Illustration error: {e}", fill='red', font=small_font)

        # Add page number
        page_num_text = str(page_num + 1)
        draw.text((500, 700), page_num_text, fill='black', font=small_font, anchor="mm")

        return img

    # def preview(self):
    #     """Launch a Gradio interface for previewing and editing the book."""
    #     def update_prompt(evt: gr.SelectData) -> str:
    #         """Update prompt when gallery image is selected."""
    #         page_num = evt.index
    #         return self._get_prompt(self.pages[page_num])
        
    #     def generate_illustration(prompt: str, evt: gr.SelectData) -> list:
    #         """Generate new illustration for selected page."""
    #         page_num = evt.index
            
    #         # Get HF token and setup API
    #         token = huggingface_hub.get_token()
    #         API_URL = f"https://api-inference.huggingface.co/models/{self.lora}"
    #         headers = {"Authorization": f"Bearer {token}"}
            
    #         try:
    #             # Generate new illustration
    #             response = requests.post(API_URL, headers=headers, json={"inputs": prompt})
    #             if response.status_code == 200:
    #                 # Save the image to a temporary file
    #                 image = Image.open(io.BytesIO(response.content))
    #                 temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    #                 image.save(temp_file.name)
                    
    #                 # Update the book's illustration
    #                 self.illustrations[page_num] = temp_file.name
                    
    #                 # Regenerate all page images
    #                 return [self._create_page_image(i) for i in range(len(self.pages))]
    #             else:
    #                 raise Exception(f"API returned status code {response.status_code}")
    #         except Exception as e:
    #             print(f"Error generating illustration: {e}")
    #             return None

    #     # Create initial page images
    #     page_images = [self._create_page_image(i) for i in range(len(self.pages))]

    #     with gr.Blocks() as demo:
    #         gr.Markdown(f"# {self.title}")
    #         if self.author:
    #             gr.Markdown(f"## by {self.author}")
            
    #         with gr.Row():
    #             prompt = gr.Textbox(
    #                 label="Illustration Prompt",
    #                 placeholder="Select a page to edit its illustration..."
    #             )
    #             generate_btn = gr.Button("Generate", variant="primary")

    #         gallery = gr.Gallery(
    #             value=page_images,
    #             columns=3,
    #             height=500,
    #             show_label=False
    #         ).style(grid=3)

    #         # Setup event handlers
    #         gallery.select(update_prompt, None, prompt)
    #         generate_btn.click(
    #             generate_illustration,
    #             inputs=[prompt],
    #             outputs=[gallery],
    #             _js="(prompt, evt) => [prompt, selected_index]",  # Pass gallery selection
    #             preprocess=False
    #         )

    #     demo.launch()